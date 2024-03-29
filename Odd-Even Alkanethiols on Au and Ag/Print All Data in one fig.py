
# -*- coding: utf-8 -*-
"""
This program is used to print the heat map based on the measurement of SAMS. It also contains the function to divide measurements in to different time 
period based on their experiment time. Still need to add methods to determine humidity and temperature. 

@author: Jack & ZZ
"""

from PIL import Image
from pptx import Presentation #pip install python-pptx
from pptx.util import Inches
import matplotlib.pyplot as plt
import pandas as pd
import numpy as np
import os
import re
from natsort import natsorted
import csv
import codecs
import urllib.request
import urllib.error
import sys
from collections import defaultdict
import seaborn as sns
from io import BytesIO



fileDnT = []  #filename that contains data and time 
morning = []  #6am - 10am
noon = []     #10am - 2pm
afternoon = []  #2pm - 6pm
evening = []   #6pm - 10pm
night=[] #10pm-6am
morning_counter = 0
noon_counter = 0
afternoon_counter = 0
evening_counter = 0
night_counter = 0
morning_total = 0
noon_total = 0
afternoon_total = 0
evening_total = 0
night_total = 0

RealFull = pd.DataFrame()
# Need to change due to different file path
pathfolder=r'E:\SMMT Group\Origin Data\New data from Martin\MAA-monoAlkylAmides\20090812(C11NHBu 3h)'

'''
Returen the experiment time from input data set
'''

def getExperienceTime(Series):
    
    
    time = []
    currentLength = 0
    
    for file in Series:
        f = open(file)
        count = 0
        for lines in f :
            
            if(lines[0].isnumeric()) :
                time.append(lines[0:-1])               
                count += 1

        
        if(count != 1):
            time[currentLength] = time[-1]

            for i in range(count):                
                if(len(time) > currentLength + 1):                   
                    del time[- 1]

                    
        currentLength += 1
        
    return time


'''
Find the related files and their correspinding times 
'''
def find_related_files(txt_files):
    related_files = []
    times = []
    for txt_file in txt_files:
        directory = os.path.dirname(txt_file)  # Get the directory of the TXT file

        base_name = os.path.basename(txt_file)  # Get the base name of the TXT file
        base_name_without_data = base_name.replace("_data.txt", "")  # Remove the "_data.txt" suffix

        files_in_directory = os.listdir(directory)
        for file in files_in_directory:
            if file.startswith(base_name_without_data) and not file.endswith(".txt"):
                related_file_path = os.path.join(directory, file)
                related_files.append(txt_file)  # Keep the original TXT file path

                # Extract time from the file content using the getExperienceTime function
                time = getExperienceTime([related_file_path])
                if time:
                    times.append(time[0])
                else:
                    times.append(None)

    # Create a DataFrame with related files and their corresponding times
    df = pd.DataFrame({"TXT File": related_files, "Time": times})
    return df

'''
Determine the time period of the measurement 
'''
def findTimePeriod(files, time):
    morning = []
    noon = []
    afternoon = []
    evening = []
    night = []
    
    for index in range(len(time)):
        file = files[index]
        if time[index][-2:] == "AM":  # Starting with AM times
            hour = int(time[index][0:time[index].index(':')])
            #print(hour)
            if 6 <= hour < 10:
                morning.append(file)
            elif hour < 6:
                night.append(file)
            elif hour == 12:
                 night.append(file)
            else:
                noon.append(file)
        else:
            hour = int(time[index][0:time[index].index(':')])
            if hour < 2:
                noon.append(file)
            elif hour == 12:
                noon.append(file)
            elif 2 <= hour < 6:
                afternoon.append(file)
            elif 6 <= hour < 10:
                evening.append(file)
            else:
                night.append(file)
    
    return morning, noon, afternoon, evening, night  

'''
Extract numbersd from file name
'''
def extract_numbers_from_filename(filename):
    numbers = re.findall(r'\d+', filename)
    return [int(num) for num in numbers]

'''
Concatenate data in a series of full pathways together with specific format in these files (voltage,absJ,J,Current,Time)
'''
def DataConcat(Series): 
    i=0
    Full=pd.DataFrame()
    dfs=pd.DataFrame()
    df=[]
    for values in Series:
        df=pd.read_csv(values,sep='\t',on_bad_lines='skip')
        if i==1:
            x=dfs.columns
            b=df.columns
            df=df.rename(columns={b[k]:x[k] for k in range(len(x))})
            Full=pd.concat([dfs,df],ignore_index=True)
        elif i>=2:
            c=df.columns
            df=df.rename(columns={c[j]:x[j] for j in range(len(c))})
            Full=pd.concat([Full,df],ignore_index=True)
        i=i+1
        dfs=df
    if i==1:
        Full=df
    new_names=['Voltage (V)','Absolute Value of J','Current Density (J)','Current','Time']
    columns=Full.columns.tolist()
    for i in range(len(new_names)):
        columns[i]=new_names[i]
    Full.columns=columns
    #if len(Full.columns)>=7:
        #Full=Full.drop(Full.columns[[5,6]],axis=1)
    #else:
        #Full=Full.drop(Full.columns[5], axis=1)
    Full=Full.iloc[:,:5]
    return Full

def add_image_slide(presentation, image_path):
    slide_layout = presentation.slide_layouts[6]  # 6 corresponds to the blank slide layout
    slide = presentation.slides.add_slide(slide_layout)
    left = Inches(1)
    top = Inches(1.5)
    pic = slide.shapes.add_picture(image_path, left, top, height=Inches(5))
    
'''
Main part
'''
grouped_files = defaultdict(list)

# Walk through subfolders and read TXT files
for root, dirs, files in os.walk(pathfolder):
    if files:
        # Extract the subfolder name from the root path
        subfolder_name = os.path.basename(root)

        for filename in files:
            if filename.endswith('.txt'):
                filepath = os.path.join(root, filename)

                # Use regular expression to extract the first number from the filename
                first_number_match = re.match(r'\D*(\d+)', filename)
                if first_number_match:
                    first_number = first_number_match.group(1)
                    grouped_files[(subfolder_name, int(first_number))].append(filepath)



# Create a dictionary to store the dataframes for each grouping
dataframes_dict = {}
my_dict=grouped_files
counter=0
countertot=[0,0,0,0,0,0,0,0,0,0]
i=0

my_dict['Extra-Ignore'] = 0
#print(my_dict)
check='sc10fc'
devcount=0
devtotal=[0,0,0,0,0,0,0,0,0,0]
exit_flag=False
fresh=[]
stdavg=np.linspace(0,9,10)
scounts=np.linspace(0,9,10)
alls=np.linspace(0,9,10)
stdcount=0
allcounts=0


presentation = Presentation()
exit_flag = False
slide_number = 0  # Slide number counter
slides_per_page = 6  # Number of plots to be shown per slide
image_files = []


left_positions = [Inches(0.5), Inches(4), Inches(7.5)]
top_positions = [Inches(1), Inches(4), Inches(7)]
image_width = Inches(3.5)
image_height = Inches(2.5)
while exit_flag == False:
    for key, nested_values in my_dict.items():
        if key[0] != check:
            # If we encounter a new carbon+date combination, save the previous set of plots in a PowerPoint
            if devtotal[3] != 0 and len(image_files) > 0:
                # Save the plots in a PowerPoint presentation
                for i in range(0, len(image_files), slides_per_page):
                    presentation = Presentation()
                    for j, image_file in enumerate(image_files[i:i+slides_per_page]):
                        left = left_positions[j % 3]
                        top = top_positions[j // 3]
                        add_image_slide(presentation, image_file, left, top, image_width, image_height)
                    pptx_filename = f"{check}_slides_{i//slides_per_page}.pptx"
                    presentation.save(pptx_filename)
                    print(f"Saved {min(slides_per_page, len(image_files)-i)} plots for {check} to {pptx_filename}")
                image_files.clear()

        if exit_flag:
            break

        # Your existing code to generate and plot the data goes here...
        # Example: (You will need to modify this part according to your actual data and plot generation)
        series_values = pd.Series(nested_values)
        try:
            Full = DataConcat(series_values)
            RealFull = pd.concat([RealFull,Full],ignore_index = True)
            
        except ValueError:
            exit_flag = True
            
for x in range(len(RealFull['Absolute Value of J'])):
    if(isinstance(RealFull['Absolute Value of J'][x],str)):
        print(RealFull['Absolute Value of J'][x])
              
Abso = np.log10(RealFull['Absolute Value of J'],where=(RealFull['Absolute Value of J']!=0))
Volts = RealFull["Voltage (V)"]
'''
plt.figure()
plt.scatter(RealFull['Voltage (V)'],Abso,color='black')
unique_x = np.unique(Volts)
mean_y = [np.mean(Abso[Volts == val]) for val in unique_x]
std_y = [np.std(Abso[Volts == val]) for val in unique_x]
color_indices = (np.arange(len(Volts)) // 43)
        
# Choose a colormap with a sufficient number of colors
cmap = plt.get_cmap('viridis')
        
     
scatter=plt.scatter(Volts, Abso, label='Data Points for Current Denisty',c=color_indices,cmap=cmap,edgecolors='black',zorder=2)
cbar = plt.colorbar()
cbar.set_ticks(np.unique(color_indices))
cbar.set_ticklabels(np.unique(color_indices) +1)
cbar.set_label('Traces')
plt.plot(unique_x, mean_y, color='r', label='Mean')
plt.fill_between(unique_x, np.array(mean_y) - np.array(std_y), np.array(mean_y) + np.array(std_y), color='gray', alpha=0.3, label='Standard Deviation')
plt.errorbar(unique_x, mean_y, yerr=std_y, fmt='none', capsize=4,color='white',zorder=1)
plt.xlabel('Voltage (V)')
plt.ylabel('Log 10 of J')
#plt.ylim(-10,6)
plt.legend(fontsize=7)
plt.grid(True)
#Currently need to changed manually
plt.title('Carbon+Date: C13 RT 3h')
#print(Abso)
        
plot_filename = f'Scatter_{key[0]}_{key[1]}.png'
plt.savefig(plot_filename)
'''        
        
fig2 = plt.figure(dpi=400)
hist = plt.hist2d(Volts, Abso, range=[[-0.5,0.5],[-10,0]], bins = [21,50],cmap=plt.cm.jet,edgecolor="black")
 
plt.xlabel('V',fontsize=15)
plt.ylabel('Log|J (A/cm\u00b2)',fontsize=15)
plt.xticks(fontsize=15)
plt.yticks(fontsize=15)
#Currently need to changed manually
#plt.title('C15', fontweight='bold', fontsize=15)
#plt.rc('grid', linestyle="-", color='black')
cbar = plt.colorbar()
cbar.ax.set_ylabel('Counts',fontsize=15)
cbar.ax.tick_params(labelsize=15)

plt.show()
allcounts += 1
'''
# Save the plot as an image in the same folder with a unique name
plot_filename = f'Heat_{key[0]}_{key[1]}.png'
plt.savefig(plot_filename)

# Add the image filename to the list
image_files.append(plot_filename)

# Continue with the rest of your code for data analysis and processing...
check=key[0]

# Save the PowerPoint presentation
for i in range(0, len(image_files), slides_per_page):
    presentation = Presentation()
    for image_file in image_files[i:i+slides_per_page]:
        add_image_slide(presentation, image_file)
    pptx_filename = f"{check}_slides_{i//slides_per_page}.pptx"
    presentation.save(pptx_filename)
    print(f"Saved {min(slides_per_page, len(image_files)-i)} plots for {check} to {pptx_filename}")
'''